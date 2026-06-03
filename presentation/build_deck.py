# -*- coding: utf-8 -*-
"""
Generates the MAIIC IFRS 9 ECL System demo / sales deck.
Dupleix Institute brand. Run: python build_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Brand palette -------------------------------------------------------
ORANGE = RGBColor(0xE8, 0x73, 0x1C)   # Dupleix accent
NAVY   = RGBColor(0x0B, 0x25, 0x45)   # Dupleix dark
GREEN  = RGBColor(0x16, 0xA3, 0x4A)   # MAIIC green
GOLD   = RGBColor(0xF5, 0x9E, 0x0B)   # MAIIC gold
RED    = RGBColor(0xD1, 0x24, 0x2F)   # MAIIC red
LIGHT  = RGBColor(0xF5, 0xF6, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x1F, 0x29, 0x37)
MUTE   = RGBColor(0x6B, 0x72, 0x80)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _fill(s, color)
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of (string, size, bold, color) OR list of paragraphs where
    each paragraph is a list of such tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if runs and not isinstance(runs[0], list):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, bold, color) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Calibri"
    return tb


def wordmark(slide, x, y, scale=1.0, light=False):
    """Text DUPLEIX INSTITUTE wordmark."""
    c1 = WHITE if light else ORANGE
    c2 = WHITE if light else NAVY
    text(slide, x, y, Inches(3.6), Inches(0.6),
         [[("DUPLEIX", int(20 * scale), True, c1),
           (" INSTITUTE", int(20 * scale), True, c2)]])
    text(slide, x, y + Inches(0.34 * scale), Inches(3.6), Inches(0.3),
         [[("Risk  |  Strategy  |  Data Analytics", int(8.5 * scale), False,
            WHITE if light else MUTE)]])


def maiic_mark(slide, x, y, scale=1.0, light=False):
    # simple leaf-ish glyph from rounded shapes + wordmark
    g = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.32 * scale), Inches(0.42 * scale))
    _fill(g, GREEN)
    g.rotation = 25
    g.shadow.inherit = False
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y + Inches(0.42 * scale),
                               Inches(0.55 * scale), Inches(0.07 * scale))
    _fill(b, GOLD)
    b.shadow.inherit = False
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y + Inches(0.50 * scale),
                               Inches(0.55 * scale), Inches(0.06 * scale))
    _fill(r, RED)
    r.shadow.inherit = False
    text(slide, x + Inches(0.7 * scale), y + Inches(0.02 * scale), Inches(3.2), Inches(0.6),
         [[("MAIIC", int(22 * scale), True, WHITE if light else GREEN)]])
    text(slide, x + Inches(0.7 * scale), y + Inches(0.36 * scale), Inches(4.2), Inches(0.3),
         [[("Malawi Agricultural & Industrial Investment Corporation plc",
            int(7.5 * scale), False, WHITE if light else MUTE)]])


PAGE = [0]


def chrome(slide, title, kicker="THE AUTOMATED IFRS 9 ECL PLATFORM"):
    rect(slide, 0, 0, SW, SH, WHITE)
    rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), SW, Inches(0.06), ORANGE)
    text(slide, Inches(0.6), Inches(0.18), Inches(9.5), Inches(0.3),
         [[(kicker, 11, True, GOLD)]])
    text(slide, Inches(0.6), Inches(0.42), Inches(11.0), Inches(0.7),
         [[(title, 26, True, WHITE)]])
    # footer
    PAGE[0] += 1
    rect(slide, 0, SH - Inches(0.32), SW, Inches(0.32), LIGHT)
    text(slide, Inches(0.6), SH - Inches(0.30), Inches(8), Inches(0.26),
         [[("Dupleix Institute  ·  Private & Confidential  ·  MAIIC IFRS 9 ECL System",
            8, False, MUTE)]])
    text(slide, SW - Inches(1.3), SH - Inches(0.30), Inches(0.8), Inches(0.26),
         [[(str(PAGE[0]), 9, True, NAVY)]], align=PP_ALIGN.RIGHT)


def bullets(slide, items, x=Inches(0.7), y=Inches(1.5), w=Inches(12.0), h=Inches(5.6),
            size=15, gap=9, sub_size=12):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap if lvl == 0 else 4)
        p.space_before = Pt(0)
        p.line_spacing = 1.05
        bullet = "●  " if lvl == 0 else "–  "
        r = p.add_run()
        r.text = bullet
        r.font.size = Pt((size if lvl == 0 else sub_size))
        r.font.bold = True
        r.font.color.rgb = GREEN if lvl == 0 else GOLD
        r.font.name = "Calibri"
        # support **bold** segments
        parts = it.split("**")
        for idx, seg in enumerate(parts):
            if seg == "":
                continue
            r = p.add_run()
            r.text = seg
            r.font.size = Pt(size if lvl == 0 else sub_size)
            r.font.bold = (idx % 2 == 1)
            r.font.color.rgb = INK if lvl == 0 else RGBColor(0x40, 0x48, 0x52)
            r.font.name = "Calibri"
        if lvl == 1:
            p.level = 1


def two_col(slide, left_title, left_items, right_title, right_items, y=Inches(1.55)):
    cw = Inches(5.95)
    for cx, ttl, items, col in [
        (Inches(0.6), left_title, left_items, GREEN),
        (Inches(6.75), right_title, right_items, ORANGE),
    ]:
        bar = rect(slide, cx, y, cw, Inches(0.5), col)
        text(slide, cx + Inches(0.2), y + Inches(0.08), cw - Inches(0.3), Inches(0.36),
             [[(ttl, 14, True, WHITE)]])
        b = slide.shapes.add_textbox(cx + Inches(0.1), y + Inches(0.65), cw - Inches(0.2), Inches(4.6))
        tf = b.text_frame
        tf.word_wrap = True
        f = True
        for it in items:
            p = tf.paragraphs[0] if f else tf.add_paragraph()
            f = False
            p.space_after = Pt(7)
            p.line_spacing = 1.05
            rr = p.add_run(); rr.text = "▸ "; rr.font.size = Pt(12); rr.font.bold = True
            rr.font.color.rgb = col; rr.font.name = "Calibri"
            for idx, seg in enumerate(it.split("**")):
                if not seg:
                    continue
                rr = p.add_run(); rr.text = seg; rr.font.size = Pt(12)
                rr.font.bold = (idx % 2 == 1); rr.font.color.rgb = INK
                rr.font.name = "Calibri"


def table(slide, headers, rows, x=Inches(0.6), y=Inches(1.5), w=Inches(12.1),
          col_w=None, fs=11, hfs=11.5):
    nrows = len(rows) + 1
    ncols = len(headers)
    h = Inches(0.45 + 0.42 * len(rows))
    gtbl = slide.shapes.add_table(nrows, ncols, x, y, w, h).table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            gtbl.columns[i].width = Emu(int(int(w) * cw / total))
    for j, htxt in enumerate(headers):
        c = gtbl.cell(0, j)
        c.text = htxt
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        pr = c.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.LEFT
        rn = pr.runs[0]; rn.font.size = Pt(hfs); rn.font.bold = True; rn.font.color.rgb = WHITE
        c.text_frame.word_wrap = True
        c.margin_left = Inches(0.06); c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gtbl.cell(i, j)
            c.text = str(val)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            pr = c.text_frame.paragraphs[0]
            rn = pr.runs[0]; rn.font.size = Pt(fs); rn.font.color.rgb = INK
            rn.font.bold = (j == 0)
            c.text_frame.word_wrap = True
            c.margin_left = Inches(0.06); c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
    return gtbl


def section_divider(title, subtitle):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, Inches(3.05), SW, Inches(0.06), ORANGE)
    text(s, Inches(0.9), Inches(2.2), Inches(11), Inches(0.4),
         [[("SECTION", 13, True, GOLD)]])
    text(s, Inches(0.9), Inches(2.45), Inches(11.5), Inches(1.0),
         [[(title, 34, True, WHITE)]])
    text(s, Inches(0.9), Inches(3.3), Inches(11), Inches(0.6),
         [[(subtitle, 15, False, RGBColor(0xC9, 0xD2, 0xDE))]])
    return s


# ============================== SLIDES ===================================

# 1. COVER
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 0, Inches(0.25), SH, ORANGE)
rect(s, Inches(0.25), 0, Inches(0.08), SH, GREEN)
wordmark(s, Inches(0.9), Inches(0.6), 1.15, light=True)
maiic_mark(s, SW - Inches(4.7), Inches(0.65), 1.0, light=True)
text(s, Inches(0.9), Inches(2.5), Inches(11.4), Inches(0.5),
     [[("PROPOSAL & SYSTEM DEMONSTRATION", 16, True, GOLD)]])
text(s, Inches(0.9), Inches(2.95), Inches(11.6), Inches(1.7),
     [[("The Automated IFRS 9", 44, True, WHITE)],
      [("Expected Credit Loss System", 44, True, WHITE)]], space_after=2)
text(s, Inches(0.9), Inches(4.85), Inches(11), Inches(0.5),
     [[("Prepared for Malawi Agricultural & Industrial Investment Corporation plc (MAIIC)",
        16, False, RGBColor(0xC9, 0xD2, 0xDE))]])
rect(s, Inches(0.9), Inches(5.55), Inches(5.6), Inches(0.04), ORANGE)
text(s, Inches(0.9), Inches(5.7), Inches(11), Inches(0.7),
     [[("Replacing the 8-hour Excel model with a secure, automated, audit-ready platform.",
        13, False, WHITE)]])
text(s, Inches(0.9), SH - Inches(0.7), Inches(11), Inches(0.4),
     [[("Dupleix Institute  ·  Risk | Strategy | Data Analytics  ·  Private & Confidential",
        11, False, MUTE)]])

# 2. AGENDA
s = prs.slides.add_slide(BLANK); chrome(s, "What we will cover today")
two_col(s, "The case for change",
        ["Who we are & our 4-year journey with MAIIC",
         "What your independent review found",
         "The true cost of the Excel model",
         "Our solution at a glance",
         "Full IFRS 9 standard compliance"],
        "The system, in depth",
        ["The engines: ECL, PD, LGD, FLI, SICR",
         "Stress testing, reports & AI analytics",
         "Security, performance & flexibility",
         "Risk → capability mitigation matrix",
         "Implementation, investment & the team"])

# 3. ABOUT DUPLEIX
s = prs.slides.add_slide(BLANK); chrome(s, "About Dupleix Institute")
bullets(s, [
    "A specialist **Risk, Strategy and Financial Consultancy** serving financial institutions and governments across **13 African countries**.",
    "Practice areas: **IFRS 9, credit scoring, pricing, stress testing, business forecasting and capital-management** model build & automation.",
    "Certified, case-study masterclasses, FRM tuition, e-learning and boardroom risk facilitation.",
    "**We designed MAIIC's original Excel IFRS 9 model in 2021** — we know your methodology end to end.",
    "This platform is the natural, audit-recommended evolution of that work: same trusted logic, automated.",
], size=16, gap=12)

# 4. THE JOURNEY
s = prs.slides.add_slide(BLANK); chrome(s, "Our 4-year journey with MAIIC")
table(s, ["Phase", "Period", "What happened"],
      [["Foundational", "2021–2022", "Dupleix builds MAIIC's Excel IFRS 9 model; S&P transition matrices, arrears-based grading"],
       ["Internal calibration", "2023", "Migration to MAIIC's own internally-derived transition matrices"],
       ["Complexity growth", "2024–2025", "Ebanker core banking go-live; FinES & Mega Farm portfolios; Excel hits its limits"],
       ["Automation", "2025", "Independent review + this automated PHP platform — validated against the Excel model"]],
      col_w=[2.2, 2.0, 8.0], fs=12)
text(s, Inches(0.6), Inches(4.6), Inches(12), Inches(1.0),
     [[("The October 2025 independent review confirmed the automated platform produces results "
        "materially consistent with the enhanced Excel model — automation without methodological disruption.",
        13, False, INK)]])

# 5. WHAT YOUR REVIEW FOUND
s = prs.slides.add_slide(BLANK); chrome(s, "What your independent review found")
table(s, ["Severity", "Risk", "Root cause"],
      [["CRITICAL", "Spreadsheet model failure", "Excel calc >8 hours, processor failures, 500+ facilities"],
       ["CRITICAL", "Missing RBM sector codes", "Prudential reporting & IFRS 9 segmentation at risk"],
       ["CRITICAL", "Manual process dependency", "Overrides & Excel manipulation without controls"],
       ["CRITICAL", "Data continuity", "Multiple accounts per facility; ID inconsistencies"],
       ["SIGNIFICANT", "Concentration risk blind spots", "Segmentation gaps hide correlated risk"],
       ["SIGNIFICANT", "Mega Farm risk mismatch", "Corporate PD/LGD applied to retail agri loans"],
       ["SIGNIFICANT", "Collateral allocation", "Single-realisation, manual proportional allocation"],
       ["EMERGING", "FLI / key-person / data quality", "No macro overlay; reliance on individual expertise"]],
      col_w=[1.8, 3.6, 6.8], fs=10.5, hfs=11)
text(s, Inches(0.6), Inches(5.75), Inches(12), Inches(0.5),
     [[("Source: Independent Review Report – IFRS 9 ECL, period ended 31 October 2025.", 10, False, MUTE)]])

# 6. COST OF EXCEL
s = prs.slides.add_slide(BLANK); chrome(s, "The true cost of staying on Excel")
two_col(s, "Operational pain (today)",
        ["Calculation runs exceeding **8 hours**",
         "Frequent **processor failures** mid-run",
         "Cannot handle monthly transition matrices at scale",
         "Portfolio already **exceeds 500 facilities** and growing",
         "Heavy **key-person dependency**, thin documentation"],
        "Business consequences",
        ["Risk of **financial misstatement**",
         "**Audit qualification** / management-letter findings",
         "**Regulatory penalties** (RBM prudential gaps)",
         "Reconciliation effort that recurs every month",
         "Growth (Mega Farm) **stalled** by tooling limits"])

section_divider("The Solution", "One automated, secure, audit-ready IFRS 9 platform — built on your own methodology.")

# 8. SOLUTION AT A GLANCE
s = prs.slides.add_slide(BLANK); chrome(s, "Our solution at a glance")
bullets(s, [
    "A **database-driven, web-based** IFRS 9 ECL platform — no spreadsheets in the critical path.",
    "Implements the **exact methodology Dupleix built for MAIIC**, validated against the Excel model.",
    "**Set-based SQL engine**: full-book ECL in **seconds**, not 8+ hours.",
    "End-to-end workflow: Portfolio setup → Loan data → PD / LGD / FLI models → ECL → Reports.",
    "Built-in **audit trail, period locking, role security, and an in-system user manual**.",
    "**Flexible & dynamic**: add portfolios, transition profiles and settings at any time — no rebuild.",
], size=15.5, gap=11)

# 9. RISK -> CAPABILITY (their own table)
s = prs.slides.add_slide(BLANK); chrome(s, "Every risk you identified — solved")
table(s, ["Risk identified in your review", "How our system mitigates it"],
      [["Spreadsheet model failure", "Automated, transactional database calculations — Excel eliminated"],
       ["Missing RBM sector codes", "Mandatory sector tagging + Data-Quality report hard-flags gaps"],
       ["Manual process dependency", "Systematic workflows, period locking & full audit trail"],
       ["Data continuity / IDs", "Contract-identity mapping & facility-level aggregation"],
       ["Collateral allocation", "Automated multi-collateral proportional discounting"],
       ["Concentration risk", "Portfolio, sector, grade & cooperative concentration reports (HHI)"],
       ["Mega Farm mismatch", "Dedicated retail-agri segmentation & credit-enhancement LGD"],
       ["FLI gap", "Built-in Forward-Looking engine (macro → regression → PD)"]],
      col_w=[5.0, 7.2], fs=11)
text(s, Inches(0.6), Inches(5.95), Inches(12), Inches(0.4),
     [[("This mapping is taken directly from your review's own Risk Mitigation Pathways.", 10, False, MUTE)]])

# 10. IFRS 9 COMPLIANCE
s = prs.slides.add_slide(BLANK); chrome(s, "Full IFRS 9 standard compliance")
two_col(s, "IFRS 9 (Section 5.5 impairment)",
        ["Three-stage model: **Stage 1 / 2 / 3**",
         "**12-month vs lifetime** ECL measurement",
         "**SICR** assessment (quantitative + qualitative)",
         "**ECL = EAD × PD × LGD** at account level",
         "Forward-looking information incorporated"],
        "Disclosure & regulatory",
        ["IFRS 7 **disclosure note tables**",
         "Opening→closing **ECL reconciliation**",
         "Stage **migration & charge/release**",
         "**RBM prudential** 5-category classification",
         "IFRS 9 ↔ RBM mapping & provision comparison"])

section_divider("Inside the System", "The engines that do the work — and how each one helps MAIIC.")

# 12. WORKFLOW
s = prs.slides.add_slide(BLANK); chrome(s, "System workflow")
steps = ["Portfolio\nSetup", "Customer &\nLoan Data", "PD Model\n(Transition)", "LGD\nModel",
         "Forward-\nLooking", "ECL\nCalculation", "Reports &\nStress Test"]
n = len(steps)
bw = Inches(1.62); gap = Inches(0.18); x0 = Inches(0.55); y0 = Inches(2.6)
cols = [GREEN, GREEN, GOLD, GOLD, ORANGE, RED, NAVY]
for i, st in enumerate(steps):
    bx = x0 + i * (bw + gap)
    box = slide_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y0, bw, Inches(1.5))
    _fill(box, cols[i]); box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = st; r.font.size = Pt(12.5)
    r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
    if i < n - 1:
        a = s.shapes.add_shape(MSO_SHAPE.CHEVRON, bx + bw + Inches(0.005), y0 + Inches(0.5),
                               Inches(0.17), Inches(0.5))
        _fill(a, RGBColor(0xB0, 0xB8, 0xC2)); a.shadow.inherit = False
text(s, Inches(0.6), Inches(4.5), Inches(12), Inches(1.2),
     [[("Each step feeds the next and is period-scoped. Models are built, validated, "
        "locked (immutable), then applied to the loan book — exactly mirroring the "
        "audit-tested Excel process, now automated and traceable.", 13, False, INK)]])

# 13. ECL ENGINE
s = prs.slides.add_slide(BLANK); chrome(s, "Engine 1 — The IFRS 9 ECL engine")
bullets(s, [
    "Computes **ECL = EAD × PD × LGD** at the account level, aggregated by stage & portfolio.",
    "**EAD** = carrying amount + undrawn commitments × utilisation (CCF).",
    "Runs inside a **single database transaction with a per-scope lock** — no two users can corrupt a run; partial failures roll back cleanly.",
    "Selectable **PD basis** (pre-FLI / post-FLI) and **LGD basis** (customer / collection / both).",
    "**Golden-number regression tests** prove the maths and protect every future change.",
    "Backed by a **CLI re-calculation** command for back-dated periods and audit re-runs.",
], size=14.5, gap=10)

# 14. PD ENGINE
s = prs.slides.add_slide(BLANK); chrome(s, "Engine 2 — PD: Transition profiles & matrices")
two_col(s, "Flexible by design",
        ["**Count or balance** methodology",
         "**Monthly & cumulative** transition matrices",
         "User-defined **transition profiles** (start/end tables, grading columns, default stage)",
         "Drag-and-drop **stage re-ordering**",
         "Internal-grade scale (A–G) for DFI reporting"],
        "Audit-ready",
        ["Draft → **Lock (immutable)** → apply to loan book",
         "Addresses prior **grade-averaging** audit finding",
         "Recalculate for refreshed data before lock",
         "Periods-cumulated transparency (audit trail)",
         "External (S&P) → **internal** PD evolution supported"])

# 15. LGD ENGINE
s = prs.slides.add_slide(BLANK); chrome(s, "Engine 3 — The LGD engine")
bullets(s, [
    "**System** mode: cure-rate & recovery-rate from Stage-3 cohort tracking;  **Manual** mode for expert input.",
    "**Monthly and cumulative** LGD, with period-list audit transparency.",
    "**Proportional multi-collateral allocation** with forced-sale discounting — directly resolves your review's collateral-allocation finding.",
    "**Agri credit-enhancement model**: off-take / warehouse-receipt / group-cooperative guarantee / AIP backing — realistic for smallholder lending, not just real estate.",
    "Lock → apply to loan book, fully traceable.",
], size=14.5, gap=11)

# 16. FLI ENGINE
s = prs.slides.add_slide(BLANK); chrome(s, "Engine 4 — The Forward-Looking (FLI) engine")
bullets(s, [
    "Closes the **macroeconomic overlay gap** flagged in your review.",
    "**Macro elements** library + **scenario profiles** (base / upside / downside) with probability weights.",
    "**Weighted forecast** and a **regression layer** (trained model or manual slope/intercept) translating a macro shock into a PD adjustment.",
    "Produces **post-FLI PD** on the loan book; ECL can be run pre- or post-FLI.",
    "Designed for Malawi agri drivers — rainfall/drought, FX, input prices — as standard stress levers.",
], size=14.5, gap=11)

# 17. SICR + STRESS
s = prs.slides.add_slide(BLANK); chrome(s, "Engines 5 & 6 — Staging/SICR & Stress Testing")
two_col(s, "Staging & SICR engine",
        ["Quantitative **thresholds** + qualitative **SICR triggers**",
         "Consistent, rule-based Stage 1/2/3 allocation",
         "Removes the **manual arrears-override** audit risk",
         "Stage migration fully reported & reconciled"],
        "Dedicated Stress-Testing module",
        ["Per-stage **PD multipliers & LGD add-ons**",
         "Agri presets: **drought, FX/input shock**",
         "Base vs stressed ECL by stage & portfolio",
         "**Save / reload scenarios**; cooperative contagion view"])

# 18. REPORTS
s = prs.slides.add_slide(BLANK); chrome(s, "The reports suite — 24+ regulatory & management reports")
table(s, ["Category", "Reports", "How it helps MAIIC"],
      [["Core ECL", "Executive Summary, ECL by Stage / Portfolio / Sector / Product / Internal Grade, Account-level trail, Portfolio trend", "Board & finance reporting; DFI internal-grade view"],
       ["Staging & Movement", "SICR trigger, Stage migration, Opening→Closing ECL recon, ECL charge/release", "Explains the P&L impairment movement"],
       ["Model Components", "PD, LGD & Collateral, EAD, Credit Risk Mitigation (agri)", "Model transparency & validation"],
       ["Forward-Looking", "Macro scenario, Scenario-weighted ECL", "Demonstrates FLI compliance"],
       ["RBM Prudential", "RBM classification, IFRS 9 vs RBM, NPL & arrears, Provision comparison, Concentration, Cooperative linkage", "Regulatory submission & concentration oversight"],
       ["Disclosure & Audit", "FS disclosure note tables, Audit & Data-Quality", "Audit-ready; hard-flags data gaps"]],
      col_w=[2.4, 6.0, 3.8], fs=10, hfs=11)

# 19. AI / ANALYTICS
s = prs.slides.add_slide(BLANK); chrome(s, "AI & analytics")
bullets(s, [
    "**AI Executive Commentary** — auto-generated narrative on the ECL position for board packs.",
    "**Early Warning System** — forward risk signals & watchlist before default.",
    "**Sensitivity Analysis** — interactive: shock PD/LGD directly, or run a macro shock through regression → PD → ECL.",
    "Every report exports to **CSV and PDF**; all tables paginate and are drill-ready.",
    "AI turns numbers into decisions — faster board reporting, earlier risk action.",
], size=15, gap=11)

# 20. SEGMENTATION
s = prs.slides.add_slide(BLANK); chrome(s, "Real portfolio segmentation for MAIIC")
bullets(s, [
    "Native segmentation for **MAIIC core, FinES (concessional), and Mega Farm retail-agri**.",
    "Derived agri portfolios: **Agri-Inputs, Farm Equipment, Irrigation, Agri Working Capital, Industrial**.",
    "PD, LGD and ECL computed **per portfolio & per sector** — no blended assumptions.",
    "Directly answers your review's **Mega Farm mismatch** and **concentration blind-spot** findings.",
    "RBM sector tagging enforced; **NPL by DPD reconciles exactly to Stage 3**.",
], size=15, gap=11)

# 21. DATA INTEGRITY
s = prs.slides.add_slide(BLANK); chrome(s, "Data integrity & governance — built in")
two_col(s, "Controls in the system",
        ["Full **audit trail** on critical changes",
         "**Period locking** → immutable closed periods",
         "**Imports activity log**: rows, exceptions, timing",
         "**Audit & Data-Quality report** flags missing sector / unmapped portfolio"],
        "Resolves review findings",
        ["Mandatory-field validation rules",
         "Facility-level aggregation (no count inflation)",
         "Contract-identity mapping across periods",
         "Reconciliation surfaced, not hidden in Excel"])

# 22. WORKSPACE
s = prs.slides.add_slide(BLANK); chrome(s, "Workspace & collaboration")
bullets(s, [
    "**Period-close checklist**: import → segmentation → PD → LGD → FLI → ECL → reports → sign-off, with progress tracking.",
    "**Role-aware**: administrators mark steps complete (recorded with who & when); everyone else gets a live read-only view.",
    "**In-system team messaging** per reporting period — no more email threads.",
    "Deep links jump straight to the right screen for each task.",
    "Knowledge lives in the system — directly mitigating **key-person dependency**.",
], size=15, gap=11)

# 23. SECURITY
s = prs.slides.add_slide(BLANK); chrome(s, "Security")
two_col(s, "Access & identity",
        ["Authenticated login; **two-factor authentication (2FA)**",
         "**Role & permission** based access control",
         "Forced password update on first login",
         "Session controls & sign-out"],
        "Integrity & assurance",
        ["**Transactional** ECL runs (atomic, locked)",
         "**Immutable** locked periods & models",
         "Comprehensive **activity / audit logging**",
         "Self-hosted — your data stays in your environment"])

# 24. PERFORMANCE
s = prs.slides.add_slide(BLANK); chrome(s, "Performance & scalability")
bullets(s, [
    "**Set-based SQL** — full-book ECL in **seconds**, versus 8+ hours in Excel.",
    "Comfortably handles the current **5,000+ row** book and rapid growth.",
    "Concurrency-safe: per-scope locking prevents corrupt parallel runs.",
    "Queue-ready architecture for very large books and scheduled runs.",
    "No processor failures, no fragile workbooks — reliability by design.",
], size=15.5, gap=12)

# 25. FLEXIBILITY
s = prs.slides.add_slide(BLANK); chrome(s, "Flexibility & dynamism")
bullets(s, [
    "**Add a new portfolio in minutes** — it flows through PD, LGD, ECL and every report automatically.",
    "**Transition profiles are fully user-defined**: tables, grading columns, count/balance, default stage, stage order.",
    "**All settings configurable**: reporting periods, currency, organisation, manuals — no code change.",
    "**In-built manual editor**: contextual help per module, editable by admins (the green ? on every screen).",
    "New regulatory views are added as catalogue-driven reports — the system grows with MAIIC.",
], size=15, gap=11)

section_divider("Working Together", "Implementation, investment and the people behind the platform.")

# 27. ROADMAP
s = prs.slides.add_slide(BLANK); chrome(s, "Implementation roadmap")
table(s, ["Stage", "Activities", "Payment milestone"],
      [["1. Initiation & setup", "Environment install, configuration, branding, user setup", "25% on signing"],
       ["2. Data migration", "Loan-book & client load, historical mapping, validation", "25%"],
       ["3. Model alignment & UAT", "PD/LGD/FLI/ECL parallel run vs Excel; user acceptance testing", "25%"],
       ["4. Go-live & handover", "Training, documentation, sign-off, support transition", "25% on go-live"]],
      col_w=[2.8, 7.2, 2.4], fs=11.5)
text(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.9),
     [[("Typical delivery: 6–10 weeks depending on data readiness. Methodology is already "
        "validated against your Excel model, which materially de-risks the timeline.", 13, False, INK)]])

# 28. INVESTMENT
s = prs.slides.add_slide(BLANK); chrome(s, "Investment & payment options")
two_col(s, "Option A — One-off licence",
        ["**USD 30,000** once-off, perpetual licence",
         "Paid in **4 implementation milestones (25% each)**",
         "Includes setup, data migration, UAT, go-live",
         "Standard warranty & handover included",
         "Best total cost of ownership long-term"],
        "Option B — Managed subscription",
        ["**USD 4,500 / month** managed service",
         "Plus a one-off **installation fee**",
         "Hosting support, maintenance & updates included",
         "Lower upfront commitment",
         "Switch to licence at any time"])
text(s, Inches(0.6), Inches(5.7), Inches(12), Inches(0.5),
     [[("Customised to MAIIC's portfolios and the recommendations from our independent Excel-model review. "
        "Final terms confirmed in the engagement letter.", 11, False, MUTE)]])

# 29. TEAM
s = prs.slides.add_slide(BLANK); chrome(s, "The team behind the platform")
table(s, ["Name", "Role", "Focus"],
      [["Themba Mazibuko", "Principal / Lead Facilitator", "Risk strategy, IFRS 9 methodology, engagement lead"],
       ["Edward Mazibuko", "Risk & Strategy Consultant", "Credit risk modelling & client advisory"],
       ["Farisai Maburitse", "Risk & Modelling Consultant", "IFRS 9 models, validation & analytics"],
       ["Wadzanai Rombe", "Financial Engineer / Developer", "System architecture & IFRS 9 engine; MSc Financial Engineering (in progress), WorldQuant University"],
       ["Kundai Muriwo", "Data & Analytics Consultant", "Data engineering, reporting & automation"]],
      col_w=[2.6, 3.0, 6.6], fs=11)
text(s, Inches(0.6), Inches(4.85), Inches(12), Inches(0.7),
     [[("Dupleix Institute — Risk | Strategy | Data Analytics. The same team that built and "
        "reviewed your IFRS 9 methodology now delivers its automation.", 12, False, INK)]])

# 30. WHY DUPLEIX / CTA
s = prs.slides.add_slide(BLANK); chrome(s, "Why Dupleix")
bullets(s, [
    "We **built your methodology** — zero re-learning, zero methodological disruption.",
    "Our independent review **validated** this platform against your Excel model.",
    "Every risk in your review has a **named capability** in this system.",
    "Local context, IFRS 9 + RBM expertise, and a delivery team that knows MAIIC.",
    "**Next step:** confirm the option, sign the engagement letter, and we begin Stage 1.",
], size=16, gap=13)

# 31. THANK YOU
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 0, Inches(0.25), SH, ORANGE)
rect(s, Inches(0.25), 0, Inches(0.08), SH, GREEN)
wordmark(s, Inches(0.9), Inches(0.7), 1.1, light=True)
text(s, Inches(0.9), Inches(2.7), Inches(11), Inches(1.0),
     [[("Thank you.", 40, True, WHITE)]])
text(s, Inches(0.9), Inches(3.7), Inches(11), Inches(0.6),
     [[("Let's move MAIIC from spreadsheets to a secure, automated, audit-ready IFRS 9 platform.",
        16, False, RGBColor(0xC9, 0xD2, 0xDE))]])
text(s, Inches(0.9), Inches(5.0), Inches(11), Inches(1.2),
     [[("Dupleix Institute", 16, True, GOLD)],
      [("www.dupleixinstitute.com  ·  Risk | Strategy | Data Analytics", 13, False, WHITE)],
      [("Private & Confidential — prepared for MAIIC plc", 11, False, MUTE)]], space_after=4)

out = r"c:\xampp\htdocs\MAICC-IFRS9\presentation\MAIIC-IFRS9-System-Demo.pptx"
prs.save(out)
print("Saved:", out, "slides:", len(prs.slides._sldIdLst))
